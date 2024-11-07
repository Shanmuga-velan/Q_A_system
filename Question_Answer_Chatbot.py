import os
import tempfile
import streamlit as st
import google.generativeai as genai
import fitz
import docx
from pptx import Presentation
from io import BytesIO
import requests
import re
from gtts import gTTS
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores.faiss import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.faiss import FAISS
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
import re

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def perform_web_search(query):
    search_url = "https://www.googleapis.com/customsearch/v1"
    params = {
        "q": query,
        "key": os.getenv("GOOGLE_API_KEY"),
        "num": 3,
    }
    response = requests.get(search_url, params=params)
    results = response.json()
    search_results = ""
    for item in results.get("items", []):
        search_results += item["snippet"] + "\n"
    return search_results

def extract_text_from_pdf(file):
    text = ""
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    return text

def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file):
    text = ""
    try:
        presentation = Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Failed to extract text from PPTX: {e}"
    return text

def extract_text_from_txt(file):
    text = file.read().decode("utf-8")
    return text

file_type_handlers = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "pptx": extract_text_from_pptx,
    "txt": extract_text_from_txt,
}

def extract_youtube_video_id(youtube_url):
    """Extract the video ID from a YouTube URL."""
    youtube_url = youtube_url.strip()
    
    match = re.search(r"v=([a-zA-Z0-9_-]{11})", youtube_url)
    if match:
        return match.group(1)
    match = re.search(r"youtu.be/([a-zA-Z0-9_-]{11})", youtube_url)
    if match:
        return match.group(1)

    return None  

def get_youtube_transcript(video_id):
    """Fetch transcript text for a given YouTube video ID."""
    try:
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        
        transcript_text = " ".join([entry['text'] for entry in transcript_list])
        return transcript_text

    except Exception as e:
        return f"Could not fetch transcript: {str(e)}"

def get_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector(chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conv():
    prompt_template = """
    You are a seminar taker teaching from a given document. Begin each explanation with an engaging and dynamic introduction 
    such as and not exactly 'Hello newbie! Today, we are going to learn about...'. Make sure the introduction is generated on the spot and unique 
    each time. Use simple language and make the lesson interesting and easy to understand. Provide detailed and thorough explanations.
    
    Context: {context}
    Question: {question}
    
    Answer (provide an answer which is in the given pdf):
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conv()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def extract_text(file, file_type):
    handler = file_type_handlers.get(file_type)
    if handler:
        return handler(file)
    else:
        return "Unsupported file type. Please try a PDF, DOCX, PPTX, or TXT file."

def summarize_text(text):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Summarize the following text: {text}")
    return response.text

def summarize_based_on_topics(text, topics):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Summarize the following text focusing on the topics {topics}: {text}")
    return response.text

def explain_concept(concept, text):
    web_search_results = perform_web_search(concept)
    combined_text = f"Based on the document:\n{text}\n\nAnd additional information from the web:\n{web_search_results}"
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Explain the concept of {concept} based on the following information: {combined_text}")
    return response.text

def get_gemini_response(question, text):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Answer the following question based on the document: {question}. Document: {text}")
    return response.text

def clean_text_for_speech(text):
    text = re.sub(r'[.,*]', '', text)  
    text = re.sub(r'\s+', ' ', text)  
    return text.strip()

def text_to_audio(text, language):
    cleaned_text = clean_text_for_speech(text)
    tts = gTTS(cleaned_text, lang=language)
    audio_file = BytesIO()
    tts.write_to_fp(audio_file)
    audio_file.seek(0)
    return audio_file

def generate_custom_quiz(topic, text):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content(f"Generate quiz questions based on the topic '{topic}' from the following text: {text}")
    questions = response.text.strip().split("\n")
    cleaned_questions = []
    for question in questions:
        if "?" in question:
            question_clean = re.sub(r"^\d+\.\s*|\d+\s*\.\s*", "", question)
            question_clean = re.sub(r"^Question\s*\d+\s*:\s*", "", question_clean)
            cleaned_questions.append(question_clean.strip())

    return cleaned_questions[:10] 

st.set_page_config(page_title="Study Helper")
st.header("Study Helper")
uploaded_file = st.file_uploader("Upload your document", type=list(file_type_handlers.keys()))

if uploaded_file:
    file_type = uploaded_file.name.split('.')[-1].lower()
    text = extract_text(uploaded_file, file_type)
    text_chunk = get_chunks(text)
    get_vector(text_chunk)
    if "Failed" not in text:
        summary = summarize_text(text)
        st.subheader("Document Summary")
        st.write(summary)
        try:
            audio_file = text_to_audio(summary, 'en')
            st.audio(audio_file, format='audio/mp3')
        except Exception as e:
            "Translation error"
        if st.button("Generate Quiz from Summary"):
            if summary:
                quiz_questions = generate_custom_quiz(summary, text)
                st.session_state.summary_quiz_questions = quiz_questions  
                st.session_state.summary_user_answers = [""] * len(quiz_questions)
                
        if "summary_quiz_questions" in st.session_state:
            st.subheader("Summary Quiz Questions")
            for i, question in enumerate(st.session_state.summary_quiz_questions):
                st.write(f"**Question {i+1}:** {question}")
                st.session_state.summary_user_answers[i] = st.text_input(f"Your answer to question {i+1}:", key=f"summary_answer_{i}", value=st.session_state.summary_user_answers[i])

            if st.button("Submit Summary Quiz"):
                correct_answers = len([answer for answer in st.session_state.summary_user_answers if answer.strip()]) 
                total_questions = len(st.session_state.summary_quiz_questions)
                st.write(f"You answered {correct_answers} out of {total_questions} questions.")
        
topics = st.text_input("Enter topics for focused summarization:", key="topics")

if st.button("Summarize Based on Topics"):
    if topics:
        topic_summary = summarize_based_on_topics(text, topics)
        st.subheader("Topic-Based Summary")
        st.write(topic_summary)
        try:
            audio_file = text_to_audio(summary, 'en')
            st.audio(audio_file, format='audio/mp3')
        except Exception as e:
            "Translation error"
        st.session_state['topic_summary'] = topic_summary  
if "topic_summary" in st.session_state and st.button("Explain Topic Summary Concept"):
    if st.session_state["topic_summary"]:
        explanation = explain_concept(st.session_state["topic_summary"], text)
        st.subheader("Concept Explanation")
        st.write(explanation)
        st.session_state['explanation'] = explanation  

if "topic_summary" in st.session_state and st.button("Generate Topic Quiz"):
    if st.session_state["topic_summary"]:
        quiz_questions = generate_custom_quiz(st.session_state["topic_summary"], text)
        st.session_state['topic_quiz_questions'] = quiz_questions  
        st.session_state['topic_user_answers'] = [""] * len(quiz_questions) 

if "topic_quiz_questions" in st.session_state:
    st.subheader("Topic Quiz Questions")
    for i, question in enumerate(st.session_state['topic_quiz_questions']):
        st.write(f"**Question {i+1}:** {question}")
        st.session_state['topic_user_answers'][i] = st.text_input(
            f"Your answer to question {i+1}:", key=f"topic_answer_{i}", 
            value=st.session_state['topic_user_answers'][i]
        )

    if st.button("Submit Topic Quiz"):
        correct_answers = len([answer for answer in st.session_state['topic_user_answers'] if answer.strip()])
        total_questions = len(st.session_state['topic_quiz_questions'])
        st.write(f"You answered {correct_answers} out of {total_questions} questions.")

concept = st.text_input("Enter a concept to get an explanation:", key="concept")
if st.button("Explain Concept"):
    if concept:
        explanation = explain_concept(concept, text)
        st.subheader("Concept Explanation")
        st.write(explanation)
        try:
            audio_file = text_to_audio(explanation, 'en')
            st.audio(audio_file, format='audio/mp3')
        except Exception as e:
            "Translation error"
custom_topic = st.text_input("Enter a topic for custom quiz generation:", key="custom_topic")

if st.button("Generate Custom Quiz"):
    if custom_topic:
        quiz_questions = generate_custom_quiz(custom_topic, text)
        st.session_state.quiz_questions = quiz_questions  
        st.session_state.user_answers = [""] * len(quiz_questions)  
    if "quiz_questions" in st.session_state:
        st.subheader("Custom Quiz Questions")
        for i, question in enumerate(st.session_state.quiz_questions):
            st.write(f"**Question {i+1}:** {question}")
            st.session_state.user_answers[i] = st.text_input(f"Your answer to question {i+1}:", key=f"answer_{i}", value=st.session_state.user_answers[i])

        if st.button("Submit Quiz"):
            correct_answers = len([answer for answer in st.session_state.user_answers if answer.strip()]) 
            total_questions = len(st.session_state.quiz_questions)
            st.write(f"You answered {correct_answers} out of {total_questions} questions.")
question = st.text_input("Ask a question based on the document:", key="question_ask")
if st.button("Generate Answer"):
    if question:
        with st.spinner("Getting response..."):
            response = get_gemini_response(question, text)
            st.subheader("Answer")
            st.write(response)
            try:
                audio_file = text_to_audio(response, 'en')
                st.audio(audio_file, format='audio/mp3')
            except Exception as e:
                "Translation error"
youtube_url = st.text_input("Enter YouTube video link:", key="youtube_link")

if st.button("Summarize YouTube Video"):
    video_id = extract_youtube_video_id(youtube_url)
    if video_id:
        youtube_text = get_youtube_transcript(video_id)
        if "Could not fetch transcript" not in youtube_text:
            youtube_summary = summarize_text(youtube_text)
            st.subheader("YouTube Video Summary")
            st.write(youtube_summary)
            try:
                audio_file = text_to_audio(youtube_summary, 'en')
                st.audio(audio_file, format='audio/mp3')
            except Exception as e:
                st.write("Audio generation error:", e)
        else:
            st.write(youtube_text)
    else:
        st.write("Invalid YouTube URL")
question = st.text_input("Enter your question about the YouTube video:", key="youtube_question")
if st.button("Ask Question about YouTube Video"):
    if question:
        video_id = extract_youtube_video_id(youtube_url)
        if video_id:
            youtube_text = get_youtube_transcript(video_id)
            if "Could not fetch transcript" not in youtube_text:
                #st.write("Transcript content:")
                #st.write(youtube_text) 
                response = get_gemini_response(question, youtube_text)
                st.subheader("Answer")
                st.write(response)
                try:
                    audio_file = text_to_audio(response, 'en')
                    st.audio(audio_file, format='audio/mp3')
                except Exception as e:
                    st.write("Audio generation error:", e)
            else:
                st.write(youtube_text)  
        else:
            st.write("Invalid YouTube URL")
